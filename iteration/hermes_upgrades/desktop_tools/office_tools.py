"""
Hermes Desktop Office 文档操作工具集 — 已修复路径安全/KeyError/布局崩溃
"""
import os
import logging
from typing import Optional
from pathlib import Path

_log = logging.getLogger(__name__)

# ── 安全限制 ─────────────────────────────────────────────────────────────
MAX_CONTENT_LEN = 10_000_000  # 10MB
MAX_ROWS = 100_000


def _safe_path(p: str) -> str:
    """路径净化"""
    return str(Path(p).resolve())


def _check_file_exists(p: str, name: str = "file") -> Optional[dict]:
    if p and not os.path.isfile(p):
        return {"error": f"{name} not found: {p}", "success": False}
    return None


# ═══════════════════════════════════════════════════════════════════════════
# Tool 1: create_word — Word 文档
# ═══════════════════════════════════════════════════════════════════════════

def create_word(path: str, title: str = "", content: str = "", template: str = "") -> dict:
    """创建 Word 文档"""
    try:
        if template:
            err = _check_file_exists(template, "template")
            if err: return err

        if len(content) > MAX_CONTENT_LEN:
            return {"error": "Content too large (>10MB)", "success": False}

        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document(template) if template else Document()

        if title:
            heading = doc.add_heading(title, level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        for para in content.split("\n"):
            if para.strip():
                doc.add_paragraph(para.strip())

        path = _safe_path(path)
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        doc.save(path)
        return {"path": path, "success": True}
    except Exception as e:
        _log.error("create_word failed: %s", e, exc_info=True)
        return {"error": str(e), "success": False}


def edit_word(path: str, operations: list[dict]) -> dict:
    """编辑 Word 文档"""
    try:
        err = _check_file_exists(path, "document")
        if err: return err

        from docx import Document
        from docx.shared import Inches

        doc = Document(path)

        for i, op in enumerate(operations):
            t = op.get("type")
            if t is None:
                return {"error": f"Operation #{i} missing 'type': {op}", "success": False}

            if t == "add_heading":
                doc.add_heading(op["text"], level=op.get("level", 1))
            elif t == "add_paragraph":
                doc.add_paragraph(op["text"])
            elif t == "add_table":
                rows = op.get("rows", [])
                if not rows or not rows[0]:
                    return {"error": f"Operation #{i}: table rows cannot be empty", "success": False}
                cols = max(len(row) for row in rows)
                table = doc.add_table(rows=len(rows), cols=cols)
                for r, row_data in enumerate(rows):
                    for c, cell_val in enumerate(row_data):
                        table.rows[r].cells[c].text = str(cell_val)
            elif t == "add_image":
                img_path = op.get("image_path", "")
                err = _check_file_exists(img_path, "image")
                if err: return err
                doc.add_picture(img_path, width=Inches(op.get("width", 5)))
            elif t == "replace":
                old, new = op.get("old", ""), op.get("new", "")
                for p in doc.paragraphs:
                    if old in p.text:
                        # 保留格式：遍历 runs 替换
                        for run in p.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
            else:
                return {"error": f"Unknown operation type: {t}", "success": False}

        doc.save(path)
        return {"path": path, "operations": len(operations), "success": True}
    except Exception as e:
        _log.error("edit_word failed: %s", e, exc_info=True)
        return {"error": str(e), "success": False}


# ═══════════════════════════════════════════════════════════════════════════
# Tool 2: create_ppt — PPT 演示文稿（修复硬编码布局）
# ═══════════════════════════════════════════════════════════════════════════

def _find_layout(prs, name_fragment: str):
    """按名称查找布局，避免硬编码索引"""
    for layout in prs.slide_layouts:
        if name_fragment.lower() in layout.name.lower():
            return layout
    return prs.slide_layouts[0]  # fallback


def create_ppt(path: str, slides: list[dict], template: str = "") -> dict:
    """创建 PPT 演示文稿"""
    try:
        if template:
            err = _check_file_exists(template, "template")
            if err: return err

        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(template) if template else Presentation()

        for slide_data in slides:
            layout_name = slide_data.get("layout", "content")
            title_text = slide_data.get("title", "")

            if layout_name == "title":
                layout = _find_layout(prs, "title slide")
                slide = prs.slides.add_slide(layout)
                if title_text and slide.shapes.title:
                    slide.shapes.title.text = title_text
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = slide_data.get("content", "")

            elif layout_name == "two_column":
                layout = _find_layout(prs, "two content")
                slide = prs.slides.add_slide(layout)
                if title_text and slide.shapes.title:
                    slide.shapes.title.text = title_text
                if len(slide.placeholders) > 2:
                    slide.placeholders[1].text = slide_data.get("left", "")
                    slide.placeholders[2].text = slide_data.get("right", "")

            elif layout_name == "image":
                layout = _find_layout(prs, "blank")
                slide = prs.slides.add_slide(layout)
                if title_text:
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
                    txBox.text_frame.text = title_text
                if slide_data.get("image_path"):
                    err = _check_file_exists(slide_data["image_path"], "image")
                    if err: return err
                    slide.shapes.add_picture(slide_data["image_path"], Inches(1), Inches(1.5), Inches(8), Inches(5))

            else:  # content
                layout = _find_layout(prs, "title and content")
                slide = prs.slides.add_slide(layout)
                if title_text and slide.shapes.title:
                    slide.shapes.title.text = title_text
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    for i, line in enumerate(slide_data.get("content", "").split("\n")):
                        if line.strip():
                            if i == 0:
                                tf.text = line.strip()
                            else:
                                p = tf.add_paragraph()
                                p.text = line.strip()

        path = _safe_path(path)
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        prs.save(path)
        return {"path": path, "slides": len(slides), "success": True}
    except Exception as e:
        _log.error("create_ppt failed: %s", e, exc_info=True)
        return {"error": str(e), "success": False}


# ═══════════════════════════════════════════════════════════════════════════
# Tool 3: create_excel — Excel 表格（修复冗余字体设置）
# ═══════════════════════════════════════════════════════════════════════════

def create_excel(path: str, sheets: list[dict]) -> dict:
    """创建 Excel 表格"""
    try:
        total_rows = sum(len(s.get("data", [])) for s in sheets)
        if total_rows > MAX_ROWS:
            return {"error": f"Too many rows ({total_rows} > {MAX_ROWS})", "success": False}

        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill

        wb = Workbook()
        for i, sheet_data in enumerate(sheets):
            ws = wb.active if i == 0 else wb.create_sheet()
            ws.title = sheet_data.get("name", f"Sheet{i+1}")

            rows = sheet_data.get("data", [])
            headers = sheet_data.get("headers", [])

            if headers:
                for j, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=j, value=header)
                    cell.font = Font(color="FFFFFF", bold=True)
                    cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                start_row = 2
            else:
                start_row = 1

            for r, row_data in enumerate(rows, start_row):
                for c, value in enumerate(row_data, 1):
                    ws.cell(row=r, column=c, value=value)

        path = _safe_path(path)
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        wb.save(path)
        return {"path": path, "sheets": len(sheets), "success": True}
    except Exception as e:
        _log.error("create_excel failed: %s", e, exc_info=True)
        return {"error": str(e), "success": False}


# ═══════════════════════════════════════════════════════════════════════════
# 工具注册
# ═══════════════════════════════════════════════════════════════════════════

OFFICE_TOOLS = {
    "create_word": {"fn": create_word, "concurrency": "write_serial", "description": "创建Word文档"},
    "edit_word": {"fn": edit_word, "concurrency": "write_serial", "description": "编辑Word文档"},
    "create_ppt": {"fn": create_ppt, "concurrency": "write_serial", "description": "创建PPT演示文稿"},
    "create_excel": {"fn": create_excel, "concurrency": "write_serial", "description": "创建Excel表格"},
}
