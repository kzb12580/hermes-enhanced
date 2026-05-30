"""Skills system — loadable skill packs for different domains."""

from __future__ import annotations

import json
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger("hermes-backend.skills")

# Skills directory
_SKILLS_DIR = Path(__file__).parent.parent / "skills"


class Skill:
    """Represents a loadable skill."""
    
    def __init__(self, name: str, description: str, triggers: list[str], content: str):
        self.name = name
        self.description = description
        self.triggers = triggers
        self.content = content
    
    def to_context(self) -> str:
        """Convert skill to context string for system prompt."""
        return f"### Skill: {self.name}\n{self.description}\n\n{self.content}"


class SkillManager:
    """Manages available skills."""
    
    def __init__(self):
        self._skills: dict[str, Skill] = {}
        self._load_builtin_skills()
    
    def _load_builtin_skills(self):
        """Load built-in skills."""
        builtin_skills = [
            Skill(
                name="code-review",
                description="Review code for quality, security, and best practices",
                triggers=["review", "check code", "audit", "code quality"],
                content="""## Code Review Process
1. Read the code carefully
2. Check for:
   - Security vulnerabilities (SQL injection, XSS, etc.)
   - Performance issues (N+1 queries, unnecessary loops)
   - Code style (PEP 8, naming conventions)
   - Error handling (try/except, edge cases)
   - Documentation (docstrings, comments)
3. Provide specific, actionable feedback
4. Suggest improvements with code examples"""
            ),
            Skill(
                name="debugging",
                description="Systematic debugging approach for finding and fixing bugs",
                triggers=["debug", "error", "bug", "fix", "issue", "problem"],
                content="""## Debugging Process
1. **Reproduce** — Understand how to trigger the bug
2. **Isolate** — Find the smallest code that reproduces it
3. **Identify** — Determine the root cause
4. **Fix** — Implement the minimal fix
5. **Test** — Verify the fix works and doesn't break other things
6. **Document** — Explain what went wrong and why

## Common Debugging Techniques
- Add print statements or logging
- Use debugger breakpoints
- Check recent changes (git diff)
- Read error messages carefully
- Search for similar issues online"""
            ),
            Skill(
                name="data-analysis",
                description="Analyze data with pandas, numpy, and visualization",
                triggers=["analyze", "data", "csv", "excel", "statistics", "chart", "graph"],
                content="""## Data Analysis Workflow
1. **Load Data** — Read CSV/Excel/JSON with pandas
2. **Explore** — Check shape, dtypes, missing values
3. **Clean** — Handle missing values, outliers, duplicates
4. **Analyze** — Group, aggregate, correlate
5. **Visualize** — Create charts with matplotlib/seaborn
6. **Report** — Summarize findings clearly

## Common Patterns
```python
import pandas as pd
import matplotlib.pyplot as plt

# Load
df = pd.read_csv('data.csv')

# Explore
print(df.shape)
print(df.info())
print(df.describe())

# Analyze
summary = df.groupby('category').agg({'value': ['mean', 'sum', 'count']})

# Visualize
df.plot(kind='bar', x='category', y='value')
plt.savefig('chart.png')
```"""
            ),
            Skill(
                name="presentation",
                description="Create professional presentations with python-pptx",
                triggers=["ppt", "presentation", "slides", "deck"],
                content="""## Presentation Creation
1. Plan structure (8-12 slides)
2. Use consistent design
3. Keep text concise
4. Add visuals when possible

## Slide Types
- Title slide
- Agenda/Overview
- Content slides (with bullet points)
- Data slides (with charts)
- Image slides
- Conclusion/Summary

## python-pptx Template
```python
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Presentation Title"
slide.placeholders[1].text = "Subtitle"

# Content slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Points"
body = slide.placeholders[1].text_frame
body.text = "First point"
body.add_paragraph().text = "Second point"

prs.save("presentation.pptx")
```"""
            ),
            Skill(
                name="web-scraping",
                description="Extract data from websites ethically",
                triggers=["scrape", "extract", "web data", "crawl"],
                content="""## Web Scraping Ethics
1. Check robots.txt first
2. Respect rate limits
3. Don't overload servers
4. Use official APIs when available

## Workflow
1. web_search to find relevant pages
2. web_extract to get content
3. Parse and structure the data
4. Save to file if needed

## Tips
- Use specific search queries
- Extract only what you need
- Handle errors gracefully
- Cache results when possible"""
            ),
            Skill(
                name="file-organization",
                description="Organize files and directories systematically",
                triggers=["organize", "clean up", "sort files", "manage files"],
                content="""## File Organization
1. **Inventory** — List all files and their types
2. **Categorize** — Group by type, date, or project
3. **Structure** — Create logical folder hierarchy
4. **Move** — Relocate files to proper locations
5. **Verify** — Check everything is accessible

## Common Structure
```
project/
├── src/          # Source code
├── docs/         # Documentation
├── data/         # Data files
├── output/       # Generated files
├── tests/        # Test files
└── README.md     # Project overview
```"""
            ),
            Skill(
                name="research",
                description="Conduct thorough research on any topic",
                triggers=["research", "find information", "learn about", "investigate"],
                content="""## Research Process
1. **Define** — Clarify what you need to know
2. **Search** — Use multiple search queries
3. **Evaluate** — Check source reliability
4. **Synthesize** — Combine information from multiple sources
5. **Present** — Organize findings clearly

## Research Tips
- Use specific, targeted queries
- Look for authoritative sources
- Cross-reference information
- Note sources for citations
- Distinguish facts from opinions"""
            ),
            Skill(
                name="excel-processing",
                description="Process Excel files with openpyxl",
                triggers=["excel", "xlsx", "spreadsheet", "workbook"],
                content="""## Excel Processing
### Read Excel
```python
import openpyxl
wb = openpyxl.load_workbook('file.xlsx')
ws = wb.active
for row in ws.iter_rows(values_only=True):
    print(row)
```

### Write Excel
```python
wb = openpyxl.Workbook()
ws = wb.active
ws['A1'] = 'Header'
ws.append([1, 2, 3])
wb.save('output.xlsx')
```

### Common Operations
- Read/write cells
- Iterate rows/columns
- Format cells (bold, color, alignment)
- Create charts
- Apply formulas"""
            ),
        ]
        
        for skill in builtin_skills:
            self._skills[skill.name] = skill
        
        logger.info("Loaded %d builtin skills", len(builtin_skills))
    
    def get_skill(self, name: str) -> Optional[Skill]:
        """Get a skill by name."""
        return self._skills.get(name)
    
    def get_all_skills(self) -> list[Skill]:
        """Get all available skills."""
        return list(self._skills.values())
    
    def get_skills_for_query(self, query: str) -> list[Skill]:
        """Get skills that match a query."""
        query_lower = query.lower()
        matched = []
        for skill in self._skills.values():
            for trigger in skill.triggers:
                if trigger in query_lower:
                    matched.append(skill)
                    break
        return matched
    
    def get_skills_context(self, query: Optional[str] = None, active_skills: Optional[list[str]] = None) -> str:
        """Get skills context for system prompt."""
        if active_skills:
            # Filter to only active skills by name
            skills = [s for s in self._skills.values() if s.name in active_skills]
        elif query:
            skills = self.get_skills_for_query(query)
        else:
            skills = self.get_all_skills()
        
        if not skills:
            return ""
        
        contexts = [skill.to_context() for skill in skills]
        return "\n\n".join(contexts)


# Global skill manager
skill_manager = SkillManager()
